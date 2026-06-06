# coding: utf-8
"""Merged PPT: new narrative structure + our best visuals. 14 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───
DEEP_BG = RGBColor(0x0C,0x19,0x29); CARD_BG = RGBColor(0x1A,0x30,0x48)
BLUE = RGBColor(0x0E,0xA5,0xE9); PURPLE = RGBColor(0x63,0x66,0xF1)
GREEN = RGBColor(0x10,0xB9,0x81); RED = RGBColor(0xEF,0x44,0x44)
ORANGE = RGBColor(0xFA,0x8C,0x16); WHITE = RGBColor(0xFF,0xFF,0xFF)
GRAY = RGBColor(0x94,0xA3,0xB8); DARK = RGBColor(0x1E,0x29,0x3B)
LIGHT_CARD = RGBColor(0xF0,0xF5,0xF9)

SW = Inches(13.333); SH = Inches(7.5)
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

# ─── Helpers ───
def dark_slide(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP_BG; bg.line.fill.background()
    sp = bg._element; sp.getparent().remove(sp); slide.shapes._spTree.insert(2, sp)
    # bottom line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.9), Inches(11.3), Inches(0.005))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.fill.fore_color.brightness = 0.15; line.line.fill.background()

def light_slide(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    sp = bg._element; sp.getparent().remove(sp); slide.shapes._spTree.insert(2, sp)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.9), Inches(11.3), Inches(0.005))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()

def TB(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.bold = bold; r.font.name = 'Microsoft YaHei'
    return txBox

def MTB(slide, left, top, width, height, lines):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run()
        r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = 'Microsoft YaHei'
    return txBox

def CARD(slide, left, top, width, height, bg_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color or CARD_BG; s.line.fill.background()
    return s

def GLOW(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.fill.fore_color.brightness = 0.92
    sp = s._element; sp.getparent().remove(sp); slide.shapes._spTree.insert(2, sp)

def SECTION_TAG(slide, num, title, dark=True):
    c = WHITE if dark else DARK
    TB(slide, Inches(1), Inches(0.3), Inches(2), Inches(0.5), num, 14, BLUE, bold=True)
    TB(slide, Inches(1), Inches(0.9), Inches(11), Inches(0.6), title, 32, c, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s1)
GLOW(s1, Inches(8), Inches(-2), Inches(8), Inches(8), BLUE)
GLOW(s1, Inches(-3), Inches(4), Inches(6), Inches(6), PURPLE)

logo = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.67), Inches(1.3), Inches(2), Inches(2))
logo.fill.solid(); logo.fill.fore_color.rgb = BLUE; logo.line.fill.background()
TB(s1, Inches(5.67), Inches(1.9), Inches(2), Inches(0.9), '+', 48, WHITE, bold=True, align=PP_ALIGN.CENTER)

TB(s1, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8), '码医 MediCode', 52, WHITE, bold=True, align=PP_ALIGN.CENTER)
TB(s1, Inches(1), Inches(4.4), Inches(11.3), Inches(0.5), 'DRG 改革 · AI 智能编码系统', 24, BLUE, align=PP_ALIGN.CENTER)

# Tagline box
tag_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.1), Inches(4.3), Inches(0.5))
tag_box.fill.solid(); tag_box.fill.fore_color.rgb = RGBColor(0x1A,0x40,0x70); tag_box.line.fill.background()
TB(s1, Inches(4.5), Inches(5.15), Inches(4.3), Inches(0.4), '最后一块拼图', 18, BLUE, bold=True, align=PP_ALIGN.CENTER)

TB(s1, Inches(1), Inches(6.2), Inches(11.3), Inches(0.3), '码医团队 · 上海对外经贸大学 · 2026', 12, WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s2)
TB(s2, Inches(1), Inches(0.5), Inches(11), Inches(0.6), '目录', 28, WHITE, bold=True)

toc = [
    ('01', '浪潮', 'DRG 全覆盖改写游戏规则', BLUE),
    ('02', '空白', '百亿刚需市场为什么没人做', PURPLE),
    ('03', '答案', '码医——AI 编码+质控一体化', GREEN),
    ('04', '商业', '三层收入模型与竞争壁垒', ORANGE),
    ('05', '我们', '核心创始人 + AI 协作模式', RED),
]
for i, (num, title, desc, color) in enumerate(toc):
    y = Inches(1.6 + i * 1.05)
    TB(s2, Inches(1.5), y, Inches(1), Inches(0.5), num, 36, color, bold=True)
    TB(s2, Inches(2.8), y, Inches(3), Inches(0.45), title, 24, WHITE, bold=True)
    TB(s2, Inches(2.8), y+Inches(0.45), Inches(8), Inches(0.35), desc, 13, GRAY)
    # line
    if i < 4:
        ln = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), y+Inches(0.88), Inches(9), Inches(0.005))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A,0x3F,0x5F); ln.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — WAVE: DRG Reform Context
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s3)
SECTION_TAG(s3, '01  浪潮', 'DRG/DIP 全面推行正在改写每一家医院的收入逻辑')

items = [
    ('2025年底', 'DRG/DIP全国全覆盖\n编码不再是IT问题', BLUE),
    ('飞检常态化', '编码错误直接导致\n拒付+罚款数十万起', PURPLE),
    ('CHS-DRG 2.0', '628个ADRG+1236个DRG\n编码复杂度倍增', GREEN),
]
for i, (tag, desc, color) in enumerate(items):
    x = Inches(1.2 + i * 3.8); y = Inches(2.2)
    CARD(s3, x, y, Inches(3.3), Inches(3.2), CARD_BG)
    TB(s3, x+Inches(0.3), y+Inches(0.3), Inches(2.7), Inches(0.5), tag, 22, color, bold=True)
    TB(s3, x+Inches(0.3), y+Inches(1.2), Inches(2.7), Inches(1.5), desc, 15, GRAY)

# 3 pain numbers at bottom
pains = [('10万+', '全国编码员缺口'), ('15%', '人工主诊断错误率'), ('100亿+', '年编码错误医保损失')]
for i, (num, label) in enumerate(pains):
    x = Inches(2.0 + i * 3.5)
    TB(s3, x, Inches(5.8), Inches(3), Inches(0.5), num, 36, RED if i<2 else ORANGE, bold=True, align=PP_ALIGN.CENTER)
    TB(s3, x, Inches(6.3), Inches(3), Inches(0.3), label, 13, GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — GAP: 4 Deadlock Categories
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s4)
SECTION_TAG(s4, '02  空白', '需求井喷，但四个供给方各有结构性死穴')

gaps = [
    ('HIS大厂', '东软/卫宁等', '能做，不想做', '编码只是子功能\nDRG不是战略重心', BLUE),
    ('AI创业公司', '森亿/左手医生', '想做，做不动', '合规2-3年+高销售成本\n被定制化拖成项目制', PURPLE),
    ('医院信息科', '5-15人编制', '想用，不会做', '不养开发团队\n薪资招不到工程师', ORANGE),
    ('一线编码员', '最痛的需求方', '最痛，最无力', '懂编码不懂技术\n无法把痛点变产品', RED),
]
for i, (name, ex, verdict, reason, color) in enumerate(gaps):
    x = Inches(0.8 + i * 3.15); y = Inches(2.0)
    CARD(s4, x, y, Inches(2.9), Inches(4.0), CARD_BG)
    TB(s4, x+Inches(0.15), y+Inches(0.15), Inches(2.6), Inches(0.35), name, 20, color, bold=True)
    TB(s4, x+Inches(0.15), y+Inches(0.55), Inches(2.6), Inches(0.25), ex, 11, GRAY)
    # verdict
    vb = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.15), y+Inches(1.0), Inches(2.6), Inches(0.65))
    vb.fill.solid(); vb.fill.fore_color.rgb = color; vb.line.fill.background()
    TB(s4, x+Inches(0.15), y+Inches(1.05), Inches(2.6), Inches(0.55), verdict, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
    TB(s4, x+Inches(0.2), y+Inches(1.85), Inches(2.5), Inches(1.8), reason, 12, GRAY)

TB(s4, Inches(1), Inches(6.3), Inches(11.3), Inches(0.35),
   '价值空白区 = 大厂不想做 x 创业做不动 x 医院不会做 x 一线无力做', 13, GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — ANSWER: Solution Overview
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s5)
SECTION_TAG(s5, '03  答案', '码医——唯一三合一的 AI 编码+DRG+质控方案')

# Flow
steps = ['病历\n输入', 'NLP\n智能编码', 'DRG\n自动分组', '质控审核\n+费用测算']
for i, step in enumerate(steps):
    x = Inches(1.5 + i * 2.8); y = Inches(2.2)
    s = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.6))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE if i%2==0 else PURPLE; s.line.fill.background()
    TB(s5, x, y+Inches(0.3), Inches(2.2), Inches(1.0), step, 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        a = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+Inches(2.3), y+Inches(0.55), Inches(0.4), Inches(0.5))
        a.fill.solid(); a.fill.fore_color.rgb = BLUE; a.line.fill.background()

# Key differentiators
diffs = ['CHS-DRG 1.2 国家标准', 'Docker一键私有化部署', '准确率94.1% 秒级响应', 'NLP+LLM双引擎架构']
for i, d in enumerate(diffs):
    x = Inches(1.5 + i * 2.8)
    TB(s5, x, Inches(4.3), Inches(2.5), Inches(0.4), d, 15, GREEN, bold=True, align=PP_ALIGN.CENTER)

TB(s5, Inches(1), Inches(5.0), Inches(11.3), Inches(1.5),
   '一份病历进入，编码推荐 + DRG分组 + 质控报告 + 费用测算同时输出。规则引擎保证底线（47条质控规则/离线可用），LLM语义理解突破天花板（Qwen2.5本地部署/复杂多诊断场景）。', 13, GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — DUAL ENGINE ARCHITECTURE (our best visual)
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s6)
SECTION_TAG(s6, '03  答案', '双层AI引擎：规则保底线 + LLM突破天花板')

# Output layer
out_items = ['编码推荐', 'DRG 分组', '质控报告', '费用测算']
for i, item in enumerate(out_items):
    x = Inches(1.5 + i * 2.8); y = Inches(1.6)
    s = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(0.85))
    s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()
    TB(s6, x, y+Inches(0.2), Inches(2.3), Inches(0.45), item, 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
TB(s6, Inches(0.5), Inches(1.75), Inches(1), Inches(0.35), '输出层', 11, GREEN, bold=True)

# Down arrows
for i in range(4):
    a = s6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.65+i*2.8), Inches(2.5), Inches(0.25), Inches(0.35))
    a.fill.solid(); a.fill.fore_color.rgb = BLUE; a.line.fill.background()

# AI Engine - split left/right
# Left: Rules
CARD(s6, Inches(1.2), Inches(3.0), Inches(5.2), Inches(1.8), RGBColor(0x15,0x3E,0x75))
TB(s6, Inches(1.5), Inches(3.1), Inches(4.5), Inches(0.35), '规则引擎 — 保底线', 16, BLUE, bold=True)
TB(s6, Inches(1.5), Inches(3.5), Inches(4.5), Inches(0.25), '任何时候都能用，不依赖GPU/网络/AI', 10, GRAY)
for ri, item in enumerate(['NLP实体识别', '医学知识图谱', '47条质控规则', 'CHS-DRG 1.2']):
    TB(s6, Inches(1.4+ri*1.25), Inches(3.95), Inches(1.2), Inches(0.4), item, 10, WHITE, align=PP_ALIGN.CENTER)

# Right: LLM
CARD(s6, Inches(6.9), Inches(3.0), Inches(5.2), Inches(1.8), RGBColor(0x3B,0x1F,0x6E))
TB(s6, Inches(7.2), Inches(3.1), Inches(4.5), Inches(0.35), 'LLM引擎 — 突破天花板', 16, PURPLE, bold=True)
TB(s6, Inches(7.2), Inches(3.5), Inches(4.5), Inches(0.25), 'Qwen2.5本地部署，理解医学语义', 10, GRAY)
for li, item in enumerate(['语义向量检索', 'LLM推理推荐', '诊断-手术一致性', '置信度评分']):
    TB(s6, Inches(7.1+li*1.25), Inches(3.95), Inches(1.2), Inches(0.4), item, 10, WHITE, align=PP_ALIGN.CENTER)

TB(s6, Inches(6.2), Inches(3.6), Inches(0.8), Inches(0.4), '+', 24, BLUE, bold=True, align=PP_ALIGN.CENTER)

# Data layer
data_items = ['病历文本\n(.txt/.docx/.pdf)', 'ICD编码库\n(920诊断+571手术)', 'CHS-DRG分组器\n(26MDC/800+ADRG)', '医保费率表\n(权重+支付标准)']
for i, item in enumerate(data_items):
    x = Inches(1.5 + i * 2.8); y = Inches(5.2)
    s = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(1.0))
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG; s.line.color.rgb = GRAY; s.line.width = Pt(0.5)
    TB(s6, x, y+Inches(0.1), Inches(2.3), Inches(0.8), item, 10, GRAY, align=PP_ALIGN.CENTER)
TB(s6, Inches(0.5), Inches(5.5), Inches(1), Inches(0.35), '数据层', 11, GRAY, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — METRICS with context
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6]); light_slide(s7)
SECTION_TAG(s7, '03  答案', '数字不说谎', dark=False)

metrics = [
    ('94.1%', '诊断编码 Top-1 准确率', '203份真实病历，8个科室\n超越人工平均水平(85-90%)', BLUE),
    ('47条', '质控规则，6大维度', '完整性·逻辑·编码·时效\n规范表达·语义质量', GREEN),
    ('<2秒', '全流程秒级响应', '规则引擎<1秒\nLLM增强2-5秒', PURPLE),
    ('Docker', '一键私有化部署', '2核4G即可运行\n病历数据不出院内网', ORANGE),
]
for i, (num, label, desc, color) in enumerate(metrics):
    x = Inches(0.8 + i * 3.15); y = Inches(2.0)
    CARD(s7, x, y, Inches(2.9), Inches(4.0), LIGHT_CARD)
    TB(s7, x+Inches(0.1), y+Inches(0.3), Inches(2.7), Inches(0.6), num, 38, color, bold=True, align=PP_ALIGN.CENTER)
    TB(s7, x+Inches(0.1), y+Inches(1.1), Inches(2.7), Inches(0.35), label, 14, DARK, bold=True, align=PP_ALIGN.CENTER)
    TB(s7, x+Inches(0.2), y+Inches(1.8), Inches(2.5), Inches(1.8), desc, 11, GRAY, align=PP_ALIGN.CENTER)

TB(s7, Inches(1), Inches(6.4), Inches(11), Inches(0.3),
   '准确率从94.1%到97%的路线图已制定：同义词库扩编 -> LLM升级 -> LoRA微调 -> 数据飞轮。详见商业计划书。', 10, GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — MARKET POSITIONING
# ═══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s8)
SECTION_TAG(s8, '03  答案', '市场真空：高AI能力 x 低成本 = 我们的位置')

# Left: 大厂
CARD(s8, Inches(1.0), Inches(2.2), Inches(5.3), Inches(3.5), CARD_BG)
TB(s8, Inches(1.3), Inches(2.4), Inches(4.5), Inches(0.4), '大厂/创业公司方案', 18, RED, bold=True)
items_l = ['报价 50-200万/家', '服务大三甲为主', '编码/DRG 二选一', '部署3-6个月', '需配套HIS系统']
for i, item in enumerate(items_l):
    TB(s8, Inches(1.5), Inches(3.0+i*0.55), Inches(4.5), Inches(0.4), 'x  '+item, 14, GRAY)

# Right: 码医
CARD(s8, Inches(7.0), Inches(2.2), Inches(5.3), Inches(3.5), RGBColor(0x0A,0x3D,0x28))
TB(s8, Inches(7.3), Inches(2.4), Inches(4.5), Inches(0.4), '码医 MediCode', 18, GREEN, bold=True)
items_r = ['年费 8-25万/院', '1万+二级医院为主', '编码+DRG+质控 三合一', 'Docker一键，当天上线', '独立部署，不依赖HIS']
for i, item in enumerate(items_r):
    TB(s8, Inches(7.5), Inches(3.0+i*0.55), Inches(4.5), Inches(0.4), '>  '+item, 14, WHITE)

# VS in middle
v_box = s8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(3.3), Inches(1.0), Inches(1.0))
v_box.fill.solid(); v_box.fill.fore_color.rgb = BLUE; v_box.line.fill.background()
TB(s8, Inches(5.8), Inches(3.5), Inches(1.0), Inches(0.6), 'VS', 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

TB(s8, Inches(1), Inches(6.2), Inches(11.3), Inches(0.35),
   '不是跟东软抢大三甲。东软200万/家的市场让大厂去打。我们的战场是1万家二级医院——他们买不起200万的方案，但必须解决编码问题。', 12, GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6]); light_slide(s9)
SECTION_TAG(s9, '04  商业', 'ROI 超过 6:1 —— 帮医院省钱，自己赚钱', dark=False)

tiers = [
    ('基础层', 'SaaS 订阅', '年费 8-25 万/院，按床位数分级\n一家三甲年均编码损失100万+\n花15万买码医，ROI > 6:1', BLUE),
    ('增长层', '超量计费', '超出免费配额按次计费\n5-15万/年/院\n医疗SaaS 留存率 > 95%', PURPLE),
    ('增值层', '定制+培训', '定制知识库 + 驻场培训\n20万+/年\n毛利率90%+，净利率32-44%', GREEN),
]
for i, (tag, title, desc, color) in enumerate(tiers):
    y = Inches(2.2 + i * 1.6)
    # number circle
    c = s9.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y+Inches(0.1), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    TB(s9, Inches(1.5), y+Inches(0.1), Inches(0.5), Inches(0.45), str(i+1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    TB(s9, Inches(2.3), y, Inches(2), Inches(0.4), tag, 16, color, bold=True)
    TB(s9, Inches(2.3), y+Inches(0.4), Inches(2), Inches(0.3), title, 20, DARK, bold=True)
    TB(s9, Inches(5.0), y, Inches(7), Inches(1.2), desc, 14, DARK)

targets = ['初期：三四线二级医院\n(DRG压力大，编码员缺)', '中期：省会三甲医院\n(高病历量，高复杂度)', '远期：医保局/卫健委\n(区域DRG监管平台)']
for i, t in enumerate(targets):
    x = Inches(1.5 + i * 3.8)
    TB(s9, x, Inches(6.0), Inches(3.3), Inches(0.8), t, 11, GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM
# ═══════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s10)
SECTION_TAG(s10, '05  我们', '核心创始人 + AI 协作，两月搭建完整产品')

# Main person card (center-left, large)
CARD(s10, Inches(1.5), Inches(2.0), Inches(5.5), Inches(3.8))
TB(s10, Inches(2.0), Inches(2.2), Inches(4.5), Inches(0.4), '项目负责人 & 全栈开发', 18, BLUE, bold=True)
TB(s10, Inches(2.0), Inches(2.7), Inches(4.5), Inches(0.3), '郑诗东和', 14, GRAY)
TB(s10, Inches(2.0), Inches(3.2), Inches(4.5), Inches(1.5),
   '人+AI协作模式\n2月内完成90+文件全栈系统\n8页面+7组API+8数据库表\n传统3-5人团队需3-6月\n203例测试准确率94.1%', 13, WHITE)

# Advisor cards (right, smaller)
advisors = [
    ('医学顾问', '对接中', '临床医学背景\n编码规则审核与质控验证\n预计省赛前确定'),
    ('商业顾问', '对接中', '医疗SaaS背景\n商业模式打磨与资源对接\n预计省赛前确定'),
]
for i, (role, status, desc) in enumerate(advisors):
    y = Inches(2.0 + i * 2.0)
    CARD(s10, Inches(7.5), y, Inches(4.8), Inches(1.7), CARD_BG)
    TB(s10, Inches(7.8), y+Inches(0.15), Inches(2.5), Inches(0.35), role, 16, WHITE, bold=True)
    TB(s10, Inches(10.3), y+Inches(0.15), Inches(1.8), Inches(0.35), status, 12, GRAY, align=PP_ALIGN.RIGHT)
    TB(s10, Inches(7.8), y+Inches(0.6), Inches(4.2), Inches(1.0), desc, 12, GRAY)

# Bottom highlight
TB(s10, Inches(1.5), Inches(6.2), Inches(10), Inches(0.35),
   '不是PPT创业：完整可运行全栈系统 + 203例测试数据，评审可现场验证任何功能。', 13, GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — MILESTONES (compact timeline)
# ═══════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s11)
SECTION_TAG(s11, '05  我们', '从 MVP 到国赛——清晰的发展路径')

ms = [
    ('2026.05', '产品完成', '全栈系统开发完成\n94.1%准确率验证', GREEN),
    ('2026.06', '校赛路演', 'BP+PPT定稿\n联系试点意向', BLUE),
    ('2026.07-09', '省赛+试点', '1-2家医院免费试点\n收集真实使用数据', PURPLE),
    ('2026.10', '国赛冲刺', '目标全国第一\n真实试点数据佐证', ORANGE),
    ('赛后', '商业化', '成立公司+首单付费\n建立区域渠道', RED),
]
for i, (date, phase, desc, color) in enumerate(ms):
    x = Inches(0.8 + i * 2.5); y = Inches(2.5)
    # line
    if i < 4:
        line = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.8), Inches(3.7), Inches(1.9), Inches(0.03))
        line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()
    # dot
    d = s11.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.8), Inches(3.55), Inches(0.32), Inches(0.32))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()
    TB(s11, x, y, Inches(2.3), Inches(0.35), date, 14, color, bold=True, align=PP_ALIGN.CENTER)
    TB(s11, x, y+Inches(0.5), Inches(2.3), Inches(0.35), phase, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    TB(s11, x, y+Inches(1.1), Inches(2.3), Inches(1.0), desc, 11, GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — MISSION QUOTE
# ═══════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s12)
GLOW(s12, Inches(4), Inches(0), Inches(6), Inches(6), BLUE)

TB(s12, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
   '让每一份病历都准确，', 38, WHITE, align=PP_ALIGN.CENTER)
TB(s12, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
   '让每一分医保基金都花在刀刃上。', 38, WHITE, align=PP_ALIGN.CENTER)
TB(s12, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.4),
   '—— 码医 MediCode', 20, BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU + CONTACT
# ═══════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(prs.slide_layouts[6]); dark_slide(s13)
GLOW(s13, Inches(8), Inches(-2), Inches(8), Inches(8), BLUE)

TB(s13, Inches(1), Inches(2.0), Inches(11.3), Inches(0.7), '谢谢观看', 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
TB(s13, Inches(1), Inches(3.0), Inches(11.3), Inches(0.4), '码医 MediCode · 欢迎交流', 20, BLUE, align=PP_ALIGN.CENTER)
TB(s13, Inches(1), Inches(4.2), Inches(11.3), Inches(0.3), '郑诗东和', 16, WHITE, align=PP_ALIGN.CENTER)
TB(s13, Inches(1), Inches(4.7), Inches(11.3), Inches(0.3), '邮箱：1975790036@qq.com  |  电话：15800565959', 14, GRAY, align=PP_ALIGN.CENTER)
TB(s13, Inches(1), Inches(5.4), Inches(11.3), Inches(0.3), '上海对外经贸大学', 12, GRAY, align=PP_ALIGN.CENTER)

# ─── SAVE ───
output_dir = os.path.join(os.path.dirname(__file__), '..', 'output')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, '码医_MediCode_路演PPT_合并版.pptx')
prs.save(output_path)
print('Saved: ' + output_path)
print('Slides: ' + str(len(prs.slides)))
