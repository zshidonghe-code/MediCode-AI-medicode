"""Generate 码医-MediCode 路演PPT .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Color palette
BLUE = RGBColor(0x0E, 0xA5, 0xE9)
PURPLE = RGBColor(0x63, 0x66, 0xF1)
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(title_text, subtitle_text=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK)
    # Gradient accent bar at bottom
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(1.5), Inches(3.9), Inches(10.3), Inches(1))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(22)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.LEFT
    return slide


def add_section_slide(number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    num = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(2), Inches(1.5))
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = f'0{number}'
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT

    ttl = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1))
    tf2 = ttl.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(38)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT
    return slide


def add_content_slide(title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    ttl = slide.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.3), Inches(0.9))
    tf = ttl.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    y = Inches(1.8)
    for item in content_items:
        if isinstance(item, tuple):
            # (bold_text, normal_text)
            box = slide.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.7))
            tf2 = box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run_b = p2.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(18)
            run_b.font.bold = True
            run_b.font.color.rgb = DARK
            run_n = p2.add_run()
            run_n.text = item[1]
            run_n.font.size = Pt(18)
            run_n.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        elif item == '---':
            y += Inches(0.2)
            continue
        else:
            box = slide.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.6))
            tf2 = box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = f'  {item}'
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        y += Inches(0.7)
    return slide


def add_stat_slide(title, stats):
    """stats: list of (number, label)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    ttl = slide.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.3), Inches(0.9))
    tf = ttl.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    cols = len(stats)
    w = Inches(11.3 / cols)
    for i, (num, label) in enumerate(stats):
        x = Inches(1) + w * i
        num_box = slide.shapes.add_textbox(x, Inches(2.8), w, Inches(1.2))
        tf2 = num_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = str(num)
        p2.font.size = Pt(48)
        p2.font.bold = True
        p2.font.color.rgb = BLUE
        p2.alignment = PP_ALIGN.CENTER

        lbl_box = slide.shapes.add_textbox(x, Inches(4.2), w, Inches(0.8))
        tf3 = lbl_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = label
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

        if i < cols - 1:
            line = slide.shapes.add_shape(1, x + w, Inches(2.5), Inches(0.02), Inches(3.5))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            line.line.fill.background()
    return slide


# ═══════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

title = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9), Inches(2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = '码医 MediCode'
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(1))
tf2 = sub.text_frame
p2 = tf2.paragraphs[0]
p2.text = 'AI医疗DRG编码与病历质控系统'
p2.font.size = Pt(28)
p2.font.color.rgb = BLUE
p2.alignment = PP_ALIGN.CENTER

sub2 = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
tf3 = sub2.text_frame
p3 = tf3.paragraphs[0]
p3.text = '编码 + DRG分组 + 质控 — 三合一'
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER

info = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(1))
tf4 = info.text_frame
p4 = tf4.paragraphs[0]
p4.text = '郑诗东和 · 上海对外经贸大学'
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
p4.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════
# SLIDE 2: Problem
# ═══════════════════════════════════════════════════

add_title_slide('每份病历都关系着一家医院的生死线',
                '2025年底 DRG/DIP 全面覆盖 — 编码决定收入')

# ═══════════════════════════════════════════════════
# SLIDE 3: Pain Points
# ═══════════════════════════════════════════════════

add_content_slide('医院面临的四大痛点', [
    ('10万+ ', '— 全国编码员缺口，三甲医院需10人实有3-5人'),
    ('10-15% ', '— 人工编码错误率，每错一个编码医院损失数千到数万元'),
    ('<10% ', '— 传统质控抽查覆盖率，大量问题病历成飞检定时炸弹'),
    ('数百亿 ', '— 每年因编码错误导致的医保基金浪费'),
    '---',
    ('', '不是医院不想做对，而是缺少高效的工具。'),
])

# ═══════════════════════════════════════════════════
# SLIDE 4: Solution
# ═══════════════════════════════════════════════════

add_content_slide('码医：AI驱动的一体化解决方案', [
    ('NLP智能编码 ', '→ ICD-10诊断 + ICD-9-CM-3手术，300+编码库，秒级输出'),
    ('DRG自动分组 ', '→ CHS-DRG 1.2方案，26 MDC，权重+费率+预估支付'),
    ('病历质控 ', '→ 100+规则 + LLM语义检查，全量100%覆盖'),
    ('数据驾驶舱 ', '→ CMI分析、收入趋势、质控排行，管理者实时决策'),
    ('智能流水线 ', '→ 一份病历进去，编码+DRG+质控+费用全流程输出'),
    '---',
    ('三合一 = ', '编码 + DRG分组 + 质控，市场上没有同类一体化产品'),
])

# ═══════════════════════════════════════════════════
# SLIDE 5: Demo Preview (Pipeline)
# ═══════════════════════════════════════════════════

add_section_slide(1, '产品演示：智能编码流水线')

# ═══════════════════════════════════════════════════
# SLIDE 6: How it works
# ═══════════════════════════════════════════════════

add_content_slide('Pipeline：四步完成全流程分析', [
    ('Step 1 | NLP智能编码 ', '— 解析病历文本 → 提取诊断和手术实体 → AI推荐ICD编码'),
    ('Step 2 | 病历质控 ', '— 规则引擎 + LLM语义检查 → 缺陷清单 + 质控评分'),
    ('Step 3 | DRG分组 ', '— CHS-DRG 1.2规则引擎 → MDC/ADRG/DRG + 权重 + CC/MCC'),
    ('Step 4 | 费用测算 ', '— 权重 × 费率 = 预估医保支付金额'),
    '---',
    ('演示模式：', '内置3个真实病历，typewriter自动播放，4步可视化动画'),
])

# ═══════════════════════════════════════════════════
# SLIDE 7: Value - The Money
# ═══════════════════════════════════════════════════

add_content_slide('每个演示环节，都换算成钱', [
    ('漏编1个次要诊断 ', '→ DRG分组重量级变轻 → 医院损失8320元'),
    ('1家三甲医院/年 ', '→ 15万份病历 × 15%编码错误率 → 年均损失100万+'),
    ('我们的SaaS定价 ', '→ 15万/年 → ROI = 6:1 以上'),
    ('全国推广后 ', '→ 每年为医保基金节省数百亿浪费'),
    '---',
    ('评委最关心的是"这个项目能不能赚钱"，答案是——不仅自己能赚钱，还能帮客户省钱。', ''),
])

# ═══════════════════════════════════════════════════
# SLIDE 8: QC Demo
# ═══════════════════════════════════════════════════

add_section_slide(2, '产品演示：AI病历质控')

# ═══════════════════════════════════════════════════
# SLIDE 9: QC Detail
# ═══════════════════════════════════════════════════

add_content_slide('从"抽查10%漏检60%"到"全量100%覆盖"', [
    ('规则引擎 ', '— 完整性、逻辑一致性、时效性、规范表达 4大类15+规则'),
    ('LLM语义检查 ', '— 诊断与手术一致性、主要诊断正确性、漏编检测'),
    ('六级缺陷分级 ', '— CRITICAL>MAJOR>MINOR>INFO，质控评分 0-100分'),
    ('采纳/忽略操作 ', '— 编码员即时审核，状态持久化不丢失'),
    '---',
    ('传统人工抽查覆盖率<10%，质量问题病历漏检率>60%。AI全量质控把覆盖率做到100%。', ''),
])

# ═══════════════════════════════════════════════════
# SLIDE 10: Dashboard
# ═══════════════════════════════════════════════════

add_section_slide(3, '数据驾驶舱：管理者的决策工具')

# ═══════════════════════════════════════════════════
# SLIDE 11: Dashboard Detail
# ═══════════════════════════════════════════════════

add_content_slide('6大数据看板，实时运营决策', [
    ('DRG运营概览 ', '→ 总病例数、AI编码率、质控通过率、CMI均值'),
    ('科室排名 ', '→ 编码数量、准确率、CMI对比，金/银/铜牌标记'),
    ('质控趋势 ', '→ 日/周/月维度，可切换 7/30/90/180天'),
    ('AI vs 人工编码对比 ', '→ 准确率趋势图，效果一目了然'),
    ('高频缺陷分析 ', '→ 12类质控问题分布，精准定位薄弱环节'),
    ('医保收入分析 ', '→ 实际收入 vs 优化预估 vs 优化空间，12个月趋势'),
])

# ═══════════════════════════════════════════════════
# SLIDE 12: Technology
# ═══════════════════════════════════════════════════

add_section_slide(4, '技术壁垒：为什么别人做不了')

# ═══════════════════════════════════════════════════
# SLIDE 13: Tech Detail
# ═══════════════════════════════════════════════════

add_content_slide('AI三层架构 + 规则引擎兜底', [
    ('第1层 NLP实体识别 ', '→ 正则 + 知识库匹配，从病历文本中提取诊断名/手术名/药物'),
    ('第2层 语义向量检索 ', '→ TF-IDF + char n-gram(1-3) + 余弦相似度，"胸口疼"→"心绞痛"'),
    ('第3层 LLM推理推荐 ', '→ Ollama Qwen2.5，理解医学语义，处理多诊断关联'),
    ('规则引擎兜底 ', '→ LLM不可用时自动切换，核心功能不中断，数据不出院'),
    '---',
    ('比东软望海的规则引擎更智能，比森亿智能多DRG付费场景，比小SaaS覆盖全流程。', ''),
])

# ═══════════════════════════════════════════════════
# SLIDE 14: Market
# ═══════════════════════════════════════════════════

add_stat_slide('市场空间', [
    ('3,200+', '三级医院'),
    ('10,000+', '二级医院'),
    ('22亿/年', '可及市场'),
    ('45亿/年', '长期空间'),
])

# ═══════════════════════════════════════════════════
# SLIDE 15: Business Model
# ═══════════════════════════════════════════════════

add_content_slide('商业模式：SaaS + 私有化部署', [
    ('SaaS订阅 ', '→ 8-25万/年（按床位数分级），年续费率95%+'),
    ('私有化部署 ', '→ 30-80万一次性授权 + 20%年维保费'),
    ('单位经济模型 ', '→ 毛利率90% | LTV 75万 | CAC 3-5万 | LTV/CAC = 15-25x'),
    ('三年目标 ', '→ 第1年10家(170万) → 第2年50家(950万) → 第3年200家(4,300万)'),
    '---',
    ('SaaS纯软件，边际成本极低，规模效应显著。首年即盈利。', ''),
])

# ═══════════════════════════════════════════════════
# SLIDE 16: Competition
# ═══════════════════════════════════════════════════

add_content_slide('竞争优势：五条不可复制的壁垒', [
    ('1. 三合一产品形态 ', '— 编码+DRG+质控打通，市场上无同类一体化产品'),
    ('2. AI语义理解 ', '— LLM理解医学语义，不是关键词匹配，需要NLP+LLM+医学知识库'),
    ('3. 双后端离线可用 ', '— Ollama本地部署+规则兜底，纯云端方案做不到'),
    ('4. 数据飞轮 ', '— 每多一家医院，编码库和质控规则就更完善'),
    ('5. 先发优势 ', '— 窗口期12-18个月，DRG/DIP改革2025年底全面覆盖'),
])

# ═══════════════════════════════════════════════════
# SLIDE 17: Social Value
# ═══════════════════════════════════════════════════

add_content_slide('社会价值：不止是一门生意', [
    ('医保基金安全 ', '— 每年3万亿+医保支出，编码准确率直接关系基金使用效率'),
    ('医疗质量提升 ', '— 从"抽查10%"到"全量100%"质控覆盖，守护医疗安全底线'),
    ('医务人员减负 ', '— 编码从5-10分钟/份降到秒级，10万缺口用AI填补'),
    ('公平就医 ', '— 让每份病历都按真实病情获得医保支付，杜绝"靠编码赚钱"'),
])

# ═══════════════════════════════════════════════════
# SLIDE 18: Team
# ═══════════════════════════════════════════════════

add_content_slide('团队：人+AI新协作模式', [
    ('郑诗东和 ', '— 项目负责人，上海对外经贸大学，商业策划+路演答辩'),
    ('Claude (AI Agent) ', '— 技术负责人，全栈开发+AI/ML/NLP+架构设计'),
    ('医学顾问（招募中） ', '— ICD/DRG编码规则审核，医院试点对接'),
    '---',
    ('不是PPT创业——完整可运行系统，8个前端页面+7组API+17个单元测试，现场可演示。', ''),
    ('人+AI协作模式——效率远超传统纯人类小团队。', ''),
])

# ═══════════════════════════════════════════════════
# SLIDE 19: Milestones
# ═══════════════════════════════════════════════════

add_content_slide('发展里程碑', [
    ('已完成 ', '— 全栈系统开发完成（编码+DRG+QC+Dashboard+Pipeline+Admin），可现场演示'),
    ('2026.05-06 ', '— BP定稿 + PPT制作 + 演示脚本打磨 + 校赛材料准备'),
    ('2026.06-07 ', '— 校赛路演 → 根据评委反馈迭代'),
    ('2026.07-08 ', '— 1-2家医院试点 + 收集真实数据 + 编码准确率优化'),
    ('2026.08-09 ', '— 省赛路演，冲击金奖'),
    ('2026.10-11 ', '— 国赛总决赛，目标全国第一'),
    ('赛后 ', '— 成立公司，启动商业化'),
])

# ═══════════════════════════════════════════════════
# SLIDE 20: Closing
# ═══════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

ttl = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9), Inches(1.5))
tf = ttl.text_frame
p = tf.paragraphs[0]
p.text = '让每一份病历都准确'
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

ttl2 = slide.shapes.add_textbox(Inches(2), Inches(3.6), Inches(9), Inches(1))
tf2 = ttl2.text_frame
p2 = tf2.paragraphs[0]
p2.text = '让每一分医保基金都花在刀刃上'
p2.font.size = Pt(28)
p2.font.color.rgb = BLUE
p2.alignment = PP_ALIGN.CENTER

info = slide.shapes.add_textbox(Inches(2), Inches(5.3), Inches(9), Inches(1))
tf3 = info.text_frame
p3 = tf3.paragraphs[0]
p3.text = '码医 MediCode · 郑诗东和 · 谢谢各位评委'
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════

output_path = r'C:\Users\Donghe\Desktop\码医-MediCode-路演PPT.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Total slides: {len(prs.slides)}')
