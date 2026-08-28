"""Generate the 12-slide MediCode pitch deck PPT for the competition.

Usage: python scripts/build_ppt.py
Output: output/码医_MediCode_路演PPT.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────────

DEEP_BG_START = RGBColor(0x0C, 0x19, 0x29)
DEEP_BG_END = RGBColor(0x16, 0x2D, 0x50)
BRAND_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
BRAND_PURPLE = RGBColor(0x63, 0x66, 0xF1)
BRAND_GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
RED = RGBColor(0xFF, 0x4D, 0x4F)
ORANGE = RGBColor(0xFA, 0x8C, 0x16)
GREEN_CHECK = RGBColor(0x52, 0xC4, 0x1A)
CARD_BG = RGBColor(0x1A, 0x30, 0x48)
LIGHT_CARD_BG = RGBColor(0xF0, 0xF5, 0xF9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_bg(slide, color_start, color_end=None):
    """Add a gradient-ish background using a solid rect (python-pptx limitation)
    We simulate by using the start color and adding decorative shapes.
    """
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color_start
    bg.line.fill.background()
    # Send to back
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_deco_glow(slide, left, top, width, height, color, alpha=0.08):
    """Add a decorative radial glow circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Approximate transparency via brightening
    shape.fill.fore_color.brightness = 0.92
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    """Add a simple text box with one run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """Add a text box with multiple paragraphs. lines = [(text, size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_card(slide, left, top, width, height, bg_color=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or CARD_BG
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.width = Pt(0.5)
    shape.line.fill.fore_color.brightness = 0.06
    return shape


def add_bottom_line(slide):
    """Add the branding line at bottom."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.9), Inches(11.333), Inches(0.005))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_BLUE
    line.fill.fore_color.brightness = 0.15
    line.line.fill.background()


def dark_slide(slide):
    """Setup a dark impact slide."""
    add_bg(slide, DEEP_BG_START)
    add_deco_glow(slide, Inches(8), Inches(-2), Inches(8), Inches(8), BRAND_BLUE)
    add_deco_glow(slide, Inches(-3), Inches(4), Inches(6), Inches(6), BRAND_PURPLE)
    add_bottom_line(slide)
    # Page number
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 '', 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


def light_slide(slide):
    """Setup a light data slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    add_bottom_line(slide)
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 '', 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None, dark=True):
    """Add standard slide title with optional subtitle."""
    c = WHITE if dark else DARK_TEXT
    sc = LIGHT_GRAY if dark else RGBColor(0x94, 0xA3, 0xB8)
    add_text_box(slide, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
                 title, 36, c, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(0.4),
                     subtitle, 16, sc)

# ─── SLIDE 1 - Cover ─────────────────────────────────────────────────────────

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_slide(slide1)

# Logo icon square
logo_shape = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.67), Inches(1.5), Inches(2), Inches(2))
logo_shape.fill.solid()
logo_shape.fill.fore_color.rgb = BRAND_BLUE
logo_shape.line.fill.background()

add_text_box(slide1, Inches(5.67), Inches(2.1), Inches(2), Inches(0.9),
             '+', 48, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide1, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
             '码医 MediCode', 52, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1), Inches(4.6), Inches(11.3), Inches(0.5),
             'AI 驱动的 DRG 智能编码与病历质控系统', 24, BRAND_BLUE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1), Inches(5.3), Inches(11.3), Inches(0.4),
             'AI驱动 · DRG编码 · 病历质控 · 医保支付', 14, LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1), Inches(6.3), Inches(11.3), Inches(0.3),
             '码医团队 · 上海对外经贸大学', 12,
             RGBColor(0xFF, 0xFF, 0xFF),
             alignment=PP_ALIGN.CENTER)

# ─── SLIDE 2 - Pain Points ───────────────────────────────────────────────────

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide2)
add_slide_title(slide2, '中国医院的 DRG 编码困局',
                '全国10万编码员缺口，每年因编码错误导致的医保损失超百亿')

# Three pain point cards
pain_points = [
    ('10万', '全国编码员缺口', '每百张床位仅配1.2名编码员，\n远低于3-5名的合理配置'),
    ('15%', '人工编码错误率', '传统编码主要诊断错误率15%，\n手术操作漏编率达25%'),
    ('100亿+', '年医保损失金额', '编码错误 → DRG分组错误 →\n医院少收钱、医保多花钱'),
]
for i, (num, title, desc) in enumerate(pain_points):
    x = Inches(1.2 + i * 3.8)
    y = Inches(2.0)
    add_card(slide2, x, y, Inches(3.3), Inches(3.8))
    # Number
    add_text_box(slide2, x, y + Inches(0.4), Inches(3.3), Inches(0.8),
                 num, 42, RED if i < 2 else ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide2, x + Inches(0.2), y + Inches(1.5), Inches(2.9), Inches(0.4),
                 title, 20, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide2, x + Inches(0.3), y + Inches(2.2), Inches(2.7), Inches(1.2),
                 desc, 13, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide2, Inches(1), Inches(5.9), Inches(11.3), Inches(0.4),
             '编码错了，医院少收钱，医保多花钱，患者数据失真',
             14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 3 - Market Gap (价值空白区分析) ─────────────────────────────────────

slide3_gap = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide3_gap)
add_slide_title(slide3_gap, '刚需市场，为什么没有好产品？',
                '四个阵营各有结构性"死穴"——市场在等一个捅破窗户纸的人')

# Four deadlock categories as cards
gap_data = [
    ('HIS大厂', '东软/卫宁等', '能做\n不想做', '编码只是子功能\nDRG不是战略重心\n投入产出比不划算', BRAND_BLUE),
    ('AI创业公司', '森亿/左手医生等', '想做\n做不动', '合规门槛2-3年\n销售成本极高\n被定制化拖成项目制', BRAND_PURPLE),
    ('医院信息科', '5-15人编制', '想用\n不会做', '不养开发团队\n薪资招不到工程师\n失败=丢饭碗', ORANGE),
    ('一线编码员', '最真实的需求方', '最痛\n最无力', '不懂技术不懂商业\n无法把痛点\n变成产品', RED),
]

for i, (name, examples, verdict, reason, color) in enumerate(gap_data):
    x = Inches(0.8 + i * 3.15)
    y = Inches(1.8)
    # Card
    add_card(slide3_gap, x, y, Inches(2.9), Inches(4.6), CARD_BG)
    # Category name
    add_text_box(slide3_gap, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.35),
                 name, 20, color, bold=True)
    add_text_box(slide3_gap, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.25),
                 examples, 11, LIGHT_GRAY)
    # Verdict box
    v_shape = slide3_gap.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(0.7))
    v_shape.fill.solid()
    v_shape.fill.fore_color.rgb = color
    v_shape.line.fill.background()
    add_text_box(slide3_gap, x + Inches(0.15), y + Inches(1.05), Inches(2.6), Inches(0.6),
                 verdict, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Reason
    add_text_box(slide3_gap, x + Inches(0.2), y + Inches(1.9), Inches(2.5), Inches(2.2),
                 reason, 12, LIGHT_GRAY)

# Bottom hook
add_text_box(slide3_gap, Inches(1), Inches(6.7), Inches(11.3), Inches(0.35),
             '价值空白区 = 大厂不想做 × 创业做不动 × 医院不会做 × 一线无力做 -> 码医的主战场',
             13, BRAND_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 4 - Solution Overview ─────────────────────────────────────────────

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide3)
add_slide_title(slide3, '码医 — 一站式的智能编码与质控方案',
                '一份病历进入，编码推荐 + DRG 分组 + 质控报告同时输出')

# Flow diagram: 4 steps
steps = ['病历\n输入', 'NLP\n智能编码', 'DRG\n自动分组', '质控审核\n+ 费用测算']
for i, step in enumerate(steps):
    x = Inches(1.5 + i * 2.8)
    y = Inches(2.0)
    # Box
    shape = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_BLUE if i % 2 == 0 else BRAND_PURPLE
    shape.line.fill.background()
    add_text_box(slide3, x, y + Inches(0.3), Inches(2.2), Inches(1.0),
                 step, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow between
    if i < 3:
        arrow = slide3.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.3), y + Inches(0.55),
            Inches(0.4), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BRAND_BLUE
        arrow.line.fill.background()

# Bottom differentiation
points = ['编码 + 分组 + 质控 + 风险评估', 'CHS-DRG 1.2 国家标准', '规则模式可验证，响应以实测为准']
for i, p in enumerate(points):
    x = Inches(2.5 + i * 3.2)
    add_multi_text(slide3, x, Inches(4.2), Inches(2.8), Inches(1.5), [
        ('✓', 28, BRAND_GREEN, True, PP_ALIGN.CENTER),
        (p, 18, WHITE, True, PP_ALIGN.CENTER),
    ])

# ─── SLIDE 4 - Tech Architecture ─────────────────────────────────────────────

slide4 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide4)
add_slide_title(slide4, '双层 AI 引擎 — 规则 + 语义双模质控',
                '底层规则引擎保底线 + 上层LLM语义突破天花板 = 双保险架构')

# Vertical architecture: Data -> AI (dual-engine) -> Output
# ── OUTPUT LAYER (top) ──
out_items = ['编码推荐', 'DRG 分组', '质控报告', '费用测算']
for i, item in enumerate(out_items):
    x = Inches(1.5 + i * 2.8)
    y = Inches(1.6)
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRAND_GREEN; shape.line.fill.background()
    add_text_box(slide4, x, y + Inches(0.2), Inches(2.3), Inches(0.5),
                 item, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(0.5), Inches(1.8), Inches(1.0), Inches(0.4),
             '输出层', 13, BRAND_GREEN, bold=True)

# Down arrows from output to engine
for i in range(4):
    ax = Inches(2.65 + i * 2.8)
    arrow = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax, Inches(2.55), Inches(0.3), Inches(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = BRAND_BLUE; arrow.line.fill.background()

# ── AI ENGINE LAYER (middle, split into two sub-engines) ──
# Left: 规则引擎 (底线)
add_card(slide4, Inches(1.2), Inches(3.1), Inches(5.2), Inches(2.0), RGBColor(0x15, 0x3E, 0x75))
add_text_box(slide4, Inches(1.5), Inches(3.2), Inches(4.5), Inches(0.4),
             '规则引擎 — 保底线', 18, BRAND_BLUE, bold=True)
add_text_box(slide4, Inches(1.5), Inches(3.7), Inches(4.5), Inches(0.3),
             '任何时候都能用，不依赖GPU/网络/AI大模型', 11, LIGHT_GRAY)
rule_items = ['NLP实体识别', '医学知识图谱', '47条质控规则', 'CHS-DRG 1.2分组器']
for ri, item in enumerate(rule_items):
    rx = Inches(1.4 + ri * 1.25)
    add_text_box(slide4, rx, Inches(4.2), Inches(1.2), Inches(0.5),
                 item, 11, WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Right: LLM引擎 (天花板)
add_card(slide4, Inches(6.9), Inches(3.1), Inches(5.2), Inches(2.0), RGBColor(0x3B, 0x1F, 0x6E))
add_text_box(slide4, Inches(7.2), Inches(3.2), Inches(4.5), Inches(0.4),
             'LLM引擎 — 突破天花板', 18, BRAND_PURPLE, bold=True)
add_text_box(slide4, Inches(7.2), Inches(3.7), Inches(4.5), Inches(0.3),
             'Qwen2.5本地部署，理解医学语义，处理复杂多诊断', 11, LIGHT_GRAY)
llm_items = ['语义向量检索', 'LLM推理推荐', '诊断-手术一致性', '编码置信度评分']
for li, item in enumerate(llm_items):
    lx = Inches(7.1 + li * 1.25)
    add_text_box(slide4, lx, Inches(4.2), Inches(1.2), Inches(0.5),
                 item, 11, WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Vertical divider between sub-engines
add_text_box(slide4, Inches(6.2), Inches(3.8), Inches(0.8), Inches(0.5),
             '+', 28, BRAND_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ── DATA LAYER (bottom) ──
data_items = ['病历文本\n(.txt/.docx/.pdf)', 'ICD编码库\n(920诊断+571手术)', 'CHS-DRG分组器\n(26MDC/628ADRG)', '医保费率表\n(权重+支付标准)']
for i, item in enumerate(data_items):
    x = Inches(1.5 + i * 2.8); y = Inches(5.4)
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(1.1))
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG; shape.line.color.rgb = LIGHT_GRAY; shape.line.width = Pt(0.5)
    add_text_box(slide4, x, y + Inches(0.15), Inches(2.3), Inches(0.8),
                 item, 11, LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(0.5), Inches(5.7), Inches(1.0), Inches(0.4),
             '数据层', 13, LIGHT_GRAY, bold=True)

# Up arrows from data to engine
for i in range(2):
    ax = Inches(3.8 + i * 6.5)
    arrow = slide4.shapes.add_shape(MSO_SHAPE.UP_ARROW, ax, Inches(2.9), Inches(0.3), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# ─── SLIDE 5 - Product Demo Preview ──────────────────────────────────────────

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
light_slide(slide5)
add_slide_title(slide5, '智能流水线 — 从病历到 DRG 支付的全流程自动化',
                '4 步流水线，秒级完成', dark=False)

# Demo steps visual (simulate the pipeline interface)
demo_steps = ['粘贴\n病历', 'NLP\n编码', 'DRG\n分组', '费用\n测算']
for i, step in enumerate(demo_steps):
    x = Inches(1.5 + i * 2.8)
    shape = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.2), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_CARD_BG
    shape.line.color.rgb = BRAND_BLUE
    shape.line.width = Pt(1)
    add_text_box(slide5, x, Inches(2.8), Inches(2.2), Inches(1.4),
                 step, 22, BRAND_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # Step number
    circle = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.8), Inches(2.0), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_BLUE
    circle.line.fill.background()
    add_text_box(slide5, x + Inches(0.8), Inches(2.05), Inches(0.6), Inches(0.5),
                 str(i + 1), 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Demo highlight boxes
highlights = [
    ('秒级响应', 'AI 自动完成\n编码 + 质控 + DRG', Inches(1.5), Inches(5.0)),
    ('流式输出', '实时展示\n分析进度', Inches(5.5), Inches(5.0)),
    ('一键保存', '全链路数据\n入数据库', Inches(9.5), Inches(5.0)),
]
for title, desc, x, y in highlights:
    add_card(slide5, x, y, Inches(3.2), Inches(1.5), RGBColor(0xF8, 0xFA, 0xFC))
    add_text_box(slide5, x + Inches(0.2), y + Inches(0.1), Inches(2.8), Inches(0.4),
                 title, 18, BRAND_BLUE, bold=True)
    add_text_box(slide5, x + Inches(0.2), y + Inches(0.6), Inches(2.8), Inches(0.8),
                 desc, 14, DARK_TEXT)

# ─── SLIDE 6 - Coding Accuracy Data ──────────────────────────────────────────

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
light_slide(slide6)
add_slide_title(slide6, '编码验证结果：4 例规则模式测试',
                '基于4例合成病历的工程回归，不代表临床准确率或人工对比结论', dark=False)

# Big numbers - left side
metrics = [
    ('2/4', '主诊断匹配', BRAND_BLUE),
    ('2/4', '手术编码匹配', LIGHT_GRAY),
    ('669.6ms', '基准平均响应', BRAND_GREEN),
]
for i, (num, label, color) in enumerate(metrics):
    y = Inches(2.0 + i * 1.6)
    add_text_box(slide6, Inches(1.5), y, Inches(4.5), Inches(0.6),
                 num, 48, color, bold=True)
    add_text_box(slide6, Inches(1.5), y + Inches(0.7), Inches(4.5), Inches(0.3),
                 label, 14, DARK_TEXT)

# Chart area - right side (simplified bar chart)
chart_bars = [50.0, 50.0, 25.0, 50.0]
for i, val in enumerate(chart_bars[:3]):
    x = Inches(6.5 + i * 0.5)
    h = Inches(val / 100 * 3.5)
    y = Inches(5.0 - h.inches)
    bar = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Emu(int(y)), Inches(0.35), Emu(int(h)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_BLUE
    bar.line.fill.background()

add_text_box(slide6, Inches(6.5), Inches(5.5), Inches(5), Inches(0.3),
             '基于4例合成病历的规则模式运行记录', 11, LIGHT_GRAY)

# ─── SLIDE 7 - QC Capability ─────────────────────────────────────────────────

slide7 = prs.slides.add_slide(prs.slide_layouts[6])
light_slide(slide7)
add_slide_title(slide7, '内涵质控 — 覆盖 6 大维度、47 条质控规则',
                '完整性、逻辑一致性、编码一致性、时效性、规范表达、语义质量', dark=False)

dims = [
    ('完整性', '12条', '出院小结四大核心\n部分齐全性检查'),
    ('逻辑一致', '8条', '诊断与性别/手术\n部位一致性'),
    ('编码一致', '7条', 'ICD编码与诊断\n文本匹配验证'),
    ('时效性', '5条', '入院/手术记录\n24h完成检查'),
    ('规范表达', '9条', '诊断名称标准\n化与病因检查'),
    ('语义质量', '6条', 'LLM驱动的深度\n语义一致性'),
]
for i, (name, rules, desc) in enumerate(dims):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(2.0 + row * 2.3)
    add_card(slide7, x, y, Inches(3.3), Inches(2.0), RGBColor(0xF8, 0xFA, 0xFC))
    add_text_box(slide7, x + Inches(0.2), y + Inches(0.15), Inches(2.9), Inches(0.3),
                 name, 18, BRAND_BLUE, bold=True)
    add_text_box(slide7, x + Inches(0.2), y + Inches(0.5), Inches(2.9), Inches(0.3),
                 rules, 14, BRAND_PURPLE, bold=True)
    add_text_box(slide7, x + Inches(0.2), y + Inches(0.9), Inches(2.9), Inches(0.9),
                 desc, 12, DARK_TEXT)

# ─── SLIDE 8 - DRG Value ─────────────────────────────────────────────────────

slide8 = prs.slides.add_slide(prs.slide_layouts[6])
light_slide(slide8)
add_slide_title(slide8, '让每一分医保基金都花在刀刃上',
                '精准编码 = 合理支付 = 医保控费 + 医院收益双赢', dark=False)

# Flow equation
eq_parts = ['漏编 1 个\n次要诊断', '权重降低\n0.35', '医院少收\n¥8,320']
for i, (text, val) in enumerate(zip(eq_parts, ['', '', ''])):
    x = Inches(1.5 + i * 3.8)
    shape = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.0), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED if i == 0 else ORANGE if i == 1 else BRAND_GREEN
    shape.line.fill.background()
    add_text_box(slide8, x, Inches(2.3), Inches(3.0), Inches(1.2),
                 text, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 2:
        arrow = slide8.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(3.1), Inches(2.6),
            Inches(0.6), Inches(0.6))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BRAND_BLUE
        arrow.line.fill.background()

# Bottom comparison
add_text_box(slide8, Inches(1.5), Inches(4.3), Inches(10), Inches(0.4),
             '全国推广后，预计每年为医保基金节省 200 亿+',
             20, BRAND_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 9 - Competitive Advantage ─────────────────────────────────────────

slide9 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide9)
add_slide_title(slide9, '竞争格局 — 我们在哪里？',
                '高AI能力×低成本 = 无人占领的右下角  |  竞品50-200万/家，码医8-25万/年')

# Competitive comparison table headers
headers = ['功能', '码医', '东软望海', '国新健康', '森亿智能']
rows_data = [
    ['NLP智能编码', '✓ 规则模式', '✓', '✗', '△'],
    ['DRG分组', '✓ CHS1.2', '✓', '✓', '✗'],
    ['AI质控', '✓ 双模', '✗', '✗', '△'],
    ['一体化', '✓', '✗', '✗', '✗'],
    ['秒级响应', '✓', '△', '✗', '✓'],
]

# Draw table
table_top = Inches(2.0)
table_left = Inches(1.5)
col_widths = [Inches(2.4), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
row_height = Inches(0.7)

# Header row
for ci, header in enumerate(headers):
    x = table_left + sum(cw for cw in ([Inches(0)] + col_widths)[:ci+1] if isinstance(x, int))
    x = table_left
    for j in range(ci):
        x += col_widths[j]
    shape = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, table_top, col_widths[ci], row_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_BLUE
    shape.line.color.rgb = RGBColor(0x1E, 0x40, 0x70)
    shape.line.width = Pt(0.5)
    add_text_box(slide9, x, table_top + Inches(0.15), col_widths[ci], Inches(0.4),
                 header, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Data rows
for ri, row_data in enumerate(rows_data):
    y = table_top + row_height + ri * row_height
    for ci, cell in enumerate(row_data):
        x = table_left
        for j in range(ci):
            x += col_widths[j]
        is_first_col = ci == 0
        shape = slide9.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, col_widths[ci], row_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG if is_first_col else DEEP_BG_START
        shape.line.color.rgb = RGBColor(0x2A, 0x3F, 0x5F)
        shape.line.width = Pt(0.5)
        cell_color = BRAND_GREEN if '✓' in cell else RED if '✗' in cell else ORANGE if '△' in cell else WHITE
        add_text_box(slide9, x, y + Inches(0.15), col_widths[ci], Inches(0.4),
                     cell, 14, cell_color, bold=True if ci == 1 else False,
                     alignment=PP_ALIGN.CENTER)

# Positioning quadrant text below table
add_text_box(slide9, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.35),
             '定位：高AI能力 × 低成本 → 右下角无人区。大厂200万/家封顶三甲，码医8万/年打开1万家二级医院。',
             12, BRAND_GREEN, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide9, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.3),
             '2025-26 DRG全面覆盖 + AI技术成熟 + 调用成本降90% = 前所未有的时间窗口',
             11, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 10 - Business Model ───────────────────────────────────────────────

slide10 = prs.slides.add_slide(prs.slide_layouts[6])
light_slide(slide10)
add_slide_title(slide10, '清晰的商业化路径',
                'SaaS 订阅 + 按量计费 + 增值服务，多元收入模型', dark=False)

# Pyramid layers
tiers = [
    ('增值层', '定制知识库 + 驻场培训', '¥20万+/年', BRAND_GREEN),
    ('增长层', '编码超量计费（超免费配额计费）', '¥5-15万/年', BRAND_PURPLE),
    ('基础层', 'SaaS 订阅（年费 8-15 万/医院）', '¥8-15万/年', BRAND_BLUE),
]
for i, (name, desc, price, color) in enumerate(tiers):
    w = Inches(4.5 + i * 1.5)
    x = Inches(4.4 - w.inches / 2)
    y = Inches(2.5 + i * 1.3)
    shape = slide10.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide10, x + Inches(0.3), y + Inches(0.1), w - Inches(0.6), Inches(0.3),
                 name, 16, WHITE, bold=True)
    add_text_box(slide10, x + Inches(0.3), y + Inches(0.4), w - Inches(2.5), Inches(0.6),
                 desc, 13, WHITE)
    add_text_box(slide10, x + w - Inches(2.8), y + Inches(0.25), Inches(2.5), Inches(0.5),
                 price, 18, WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# Target customers
customers = [
    ('初期', '三四线城市二级医院\nDRG改革压力大、编码员短缺'),
    ('中期', '省会城市三甲医院\n高病历量、高编码复杂度'),
    ('远期', '医保局/卫健委\n区域级DRG监管平台'),
]
for i, (phase, desc) in enumerate(customers):
    x = Inches(1.5 + i * 3.8)
    add_text_box(slide10, x, Inches(5.8), Inches(3.3), Inches(0.3),
                 phase, 16, BRAND_BLUE, bold=True)
    add_text_box(slide10, x, Inches(6.2), Inches(3.3), Inches(0.8),
                 desc, 12, DARK_TEXT)

# ─── SLIDE 11 - Team & Milestones ────────────────────────────────────────────

slide11 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide11)
add_slide_title(slide11, '核心创始人 + AI协作模式',
                '一人+AI两月完成全栈开发，效率超过传统3-5人团队')

# Team members
team = [
    ('[创始人]', '项目负责人 & 全栈开发', '人+AI协作模式\n2月完成90+文件全栈系统\n效率超过传统3-5人团队'),
    ('[医疗]', '医学顾问（对接中）', '临床医学背景\n编码规则审核与质控验证\n预计省赛前确定'),
    ('[商业]', '商业顾问（对接中）', '医疗SaaS/信息化背景\n商业模式打磨与资源对接\n预计省赛前确定'),
]
for i, (avatar, role, bio) in enumerate(team):
    x = Inches(1.5 + i * 3.8)
    add_card(slide11, x, Inches(2.0), Inches(3.3), Inches(2.5))
    # Avatar placeholder
    circle = slide11.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1.15), Inches(2.2), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_BLUE
    circle.line.fill.background()
    add_text_box(slide11, x + Inches(1.15), Inches(2.35), Inches(1.0), Inches(0.7),
                 avatar[:2], 20, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide11, x + Inches(0.2), Inches(3.4), Inches(2.9), Inches(0.3),
                 role, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide11, x + Inches(0.2), Inches(3.8), Inches(2.9), Inches(0.6),
                 bio, 12, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Timeline
milestones = ['7月\nMVP完成', '8月\n试点签约', '9月\n省赛晋级', '10月\n国赛冲刺']
for i, ms in enumerate(milestones):
    x = Inches(2.5 + i * 2.5)
    y = Inches(5.2)
    # Line
    if i < 3:
        line = slide11.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.5), y + Inches(0.15),
            Inches(2.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = BRAND_BLUE
        line.line.fill.background()
    # Dot
    dot = slide11.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.5), y, Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BRAND_GREEN if i < 1 else BRAND_BLUE
    dot.line.fill.background()
    add_text_box(slide11, x + Inches(0.1), y + Inches(0.5), Inches(1.2), Inches(0.7),
                 ms, 12, WHITE, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 12 - Closing ──────────────────────────────────────────────────────

slide12 = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(slide12)

add_text_box(slide12, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
             '让每一份病历都准确', 42, WHITE, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8),
             '让每一分医保基金都花在刀刃上', 42, WHITE, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.4),
             '码医 MediCode', 24, BRAND_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.3),
             '谢谢各位评委老师', 18, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.3),
             '码医团队 · 上海对外经贸大学', 12, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─── Save ────────────────────────────────────────────────────────────────────

output_dir = os.path.join(os.path.dirname(__file__), '..', 'output')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, '码医_MediCode_路演PPT.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
